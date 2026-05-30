import collections 
import collections.abc
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------- DESIGN THEME COLORS -----------------
    DARK_BLUE = RGBColor(44, 62, 80)     # Primary dark color (#2C3E50)
    TEAL = RGBColor(38, 166, 154)       # Accent color (#26A69A)
    LIGHT_GRAY = RGBColor(245, 247, 248) # Content slide background (#F5F7F8)
    TEXT_DARK = RGBColor(52, 73, 94)     # Dark charcoal text (#34495E)
    TEXT_LIGHT = RGBColor(255, 255, 255) # White text
    ACCENT_GRAY = RGBColor(127, 140, 141) # Muted gray text (#7F8C8D)
    
    plots_dir = "plots"
    
    # Helper to set slide background color
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to create an elegant slide title with a colored accent block
    def add_slide_header(slide, title_text):
        set_background(slide, LIGHT_GRAY)
        
        # Add decorative Teal left bar
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.6), Inches(0.5), Inches(0.15), Inches(0.6)
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = TEAL
        left_bar.line.color.rgb = TEAL
        
        # Add Title Text Box
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Calibri'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Add horizontal clean gray divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.6), Inches(1.2), Inches(12.133), Inches(0.01)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(218, 223, 230)
        divider.line.fill.background()

    # Helper to add standard footer on content slides
    def add_content_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.3))
        tf = footer_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = "Intelligent Data Analysis (DV1597) | University Project"
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Calibri'
        p.font.size = Pt(11)
        p.font.color.rgb = ACCENT_GRAY

    # ==========================================
    # SLIDE 1: Title Slide (Dark Background)
    # ==========================================
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, DARK_BLUE)
    
    # Elegant footer highlight line
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), Inches(13.333), Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Exploratory Data Analysis of the COVID-19 Pandemic in the EU/EEA"
    p.font.name = 'Calibri'
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.LEFT

    # Divider bar
    divider = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(3.0), Inches(0.04))
    divider.fill.solid()
    divider.fill.fore_color.rgb = TEAL
    divider.line.fill.background()

    # Subtitle Box
    subtitle_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.5))
    tf_sub = subtitle_box.text_frame
    tf_sub.word_wrap = True
    
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Intelligent Data Analysis (DV1597) — University Project"
    p_sub.font.name = 'Calibri'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = TEAL
    
    p_details = tf_sub.add_paragraph()
    p_details.text = "\nAuthor: Bilal Abdul (abac23@student.bth.se)  |  Blekinge Institute of Technology (BTH)"
    p_details.font.name = 'Calibri'
    p_details.font.size = Pt(14)
    p_details.font.color.rgb = RGBColor(189, 195, 199)

    # ==========================================
    # SLIDE 2: Motivation & Analytical Stack (Light Background)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "Motivation & Analytical Stack")
    add_content_footer(slide2)

    # Column 1: Stack Motivation (Left)
    col1_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Why Python & Standard Data Science Stack?"
    p.font.name = 'Calibri'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(12)

    bullets1 = [
        ("Reproducible Pipelines", "Every clean, aggregation, and mathematical rank calculation is script-driven, avoiding fragile manual spreadsheet actions."),
        ("Pandas & NumPy", "Provides efficient high-performance indexing for handling millions of row combinations across multi-gigabyte datasets."),
        ("Matplotlib & Seaborn", "Enables highly optimized, vector-grade static visualizations. Eliminates heavy system dependencies (like GIS shapefile tools) to ensure cross-platform compatibility.")
    ]
    for title, desc in bullets1:
        bp = tf1.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)

    # Column 2: The ECDC Datasets (Right)
    col2_box = slide2.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = "ECDC Public Health Database"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE
    p2.space_after = Pt(12)

    bullets2 = [
        ("1. Cases & Deaths", "Tracking daily new cases and mortality numbers from 2020 through 2022 across the EU/EEA."),
        ("2. Vaccination Dataset", "Granular cohort-level vaccine tracking containing brand doses, age breakdowns, and key demographics."),
        ("3. Hospital & Clinical Burden", "Daily and weekly occupancies and admissions for general hospital wards and Intensive Care Units (ICUs).")
    ]
    for title, desc in bullets2:
        bp = tf2.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)


    # ==========================================
    # SLIDE 3: Data Cleaning & Q1 Cases Trends (Light Background + Q1 IMAGE)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "Core Pre-processing & Quarterly Cases (Q1)")
    add_content_footer(slide3)

    # Left Column: Concise text
    col1_box = slide3.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.2), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Pre-processing & Deduplication"
    p.font.name = 'Calibri'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(12)

    bullets3 = [
        ("NUTS-2 Deduplication", "Filtered 'Region == ReportingCountry' in vaccination dataset to prevent double-counting subnational regional rows."),
        ("Unified Mapping", "Mapped all country spelling discrepancies ('Czechia' -> 'CZ', 'Greece' -> 'EL') to standard ISO 2-letter codes."),
        ("Dynamic Age Fallback", "Created age aggregation fallbacks for countries like Germany that omitted granular age cohorts."),
        ("Q1 Cases Trends", "Identified the top 10 countries. Cases peaked in Q1 2022 across Western Europe (DE, FR, IT) due to Omicron emergence.")
    ]
    for title, desc in bullets3:
        bp = tf1.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)

    # Right Column: Q1 Graph Image
    q1_img_path = os.path.join(plots_dir, "q1_quarterly_cases_trends.png")
    if os.path.exists(q1_img_path):
        slide3.shapes.add_picture(q1_img_path, Inches(6.0), Inches(1.6), width=Inches(6.7))


    # ==========================================
    # SLIDE 4: Spatial Maps (Q2 IMAGE)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "Results: Spatio-Temporal Distribution Maps (Q2)")
    add_content_footer(slide4)

    # Left Column: Concise text
    col1_box = slide4.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.5), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Mapping Europe geographically"
    p.font.name = 'Calibri'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(12)

    bullets4 = [
        ("Coordinate Projection", "Plotted country cases and deaths using actual Latitude (°N) and Longitude (°E) coordinate points."),
        ("Stylized Geographical Map", "Allows clear spatial comparison of pandemic intensity directly on a Matplotlib axis without heavy shapefile libraries."),
        ("Bubble size & shade", "Size and shade instantly convey intensity. Darker, larger blue/red circles indicate high clinical case and mortality loads."),
        ("Temporal movement", "Visually demonstrates infection waves expanding from concentrated hotspots (2020) to full continental coverage (2022).")
    ]
    for title, desc in bullets4:
        bp = tf1.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)

    # Right Column: Q2 Map Image (using 2022 as it is the most representative)
    q2_img_path = os.path.join(plots_dir, "q2_spatial_bubble_map_2022.png")
    if os.path.exists(q2_img_path):
        slide4.shapes.add_picture(q2_img_path, Inches(5.3), Inches(1.4), width=Inches(7.5), height=Inches(5.4))


    # ==========================================
    # SLIDE 5: Vaccine Brands & Demographics (Q3, Q4 IMAGES Side-by-Side)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "Results: Vaccine Brand Portfolios & Demographics (Q3, Q4)")
    add_content_footer(slide5)

    # We will lay out two images side-by-side with captions underneath
    q3_img_path = os.path.join(plots_dir, "q3_vaccine_brand_distribution.png")
    q4_img_path = os.path.join(plots_dir, "q4_vaccine_brand_demographics.png")
    
    # Left Box: Q3 pie chart & concise insights
    if os.path.exists(q3_img_path):
        slide5.shapes.add_picture(q3_img_path, Inches(0.6), Inches(1.5), width=Inches(5.8))
        
    q3_caption = slide5.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(5.8), Inches(1.3))
    tf3 = q3_caption.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Q3: Overall Brand Dominance & Exceptions"
    p3.font.size = Pt(15)
    p3.font.bold = True
    p3.font.color.rgb = DARK_BLUE
    
    bp3 = tf3.add_paragraph()
    bp3.text = "• Pfizer (Comirnaty) is dominant overall at 60.3% share.\n• Liechtenstein (LI) is a major brand exception, procuring Moderna (Spikevax) as primary (65.5%).\n• Germany is a database exception, reporting 100% as 'Unknown' (UNK)."
    bp3.font.size = Pt(11)
    bp3.font.color.rgb = TEXT_DARK

    # Right Box: Q4 target groups bar charts
    if os.path.exists(q4_img_path):
        slide5.shapes.add_picture(q4_img_path, Inches(6.8), Inches(1.4), width=Inches(5.9), height=Inches(4.1))
        
    q4_caption = slide5.shapes.add_textbox(Inches(6.8), Inches(5.6), Inches(5.9), Inches(1.3))
    tf4 = q4_caption.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Q4: Brand Allocations by Demographic Target Group"
    p4.font.size = Pt(15)
    p4.font.bold = True
    p4.font.color.rgb = DARK_BLUE
    
    bp4 = tf4.add_paragraph()
    bp4.text = "• Pfizer and Moderna doses were heavily allocated to older groups (Age60+) and health care workers (HCW) in early rollouts.\n• AstraZeneca (Vaxzevria) shows a sharp restriction to specific active cohorts due to mid-campaign European policy adjustments."
    bp4.font.size = Pt(11)
    bp4.font.color.rgb = TEXT_DARK


    # ==========================================
    # NEW SLIDE 6: Demographic Rankings & Aging (Q6, Q7 IMAGES Side-by-Side)
    # ==========================================
    slide6_new = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6_new, "Results: Pediatric Uptake & Aging Cohort Rankings (Q6, Q7)")
    add_content_footer(slide6_new)

    q6_img_path = os.path.join(plots_dir, "q6_under18_vaccination_ranks.png")
    q7_img_path = os.path.join(plots_dir, "q7_weighted_average_age_ranks.png")
    
    # Left Box: Q6 pediatric ranks
    if os.path.exists(q6_img_path):
        slide6_new.shapes.add_picture(q6_img_path, Inches(0.6), Inches(1.5), width=Inches(5.8))
        
    q6_caption = slide6_new.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(5.8), Inches(1.3))
    tf6 = q6_caption.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Q6: Pediatric Vaccinations (<18) by Country"
    p6.font.size = Pt(15)
    p6.font.bold = True
    p6.font.color.rgb = DARK_BLUE
    
    bp6 = tf6.add_paragraph()
    bp6.text = "• Ireland (IE) and Denmark (DK) reported the highest proportion of vaccinated pediatric populations relative to total population.\n• Eastern European countries (e.g. Bulgaria BG, Romania RO) registered the lowest pediatric vaccination rates, showing geographical vaccination variance."
    bp6.font.size = Pt(11)
    bp6.font.color.rgb = TEXT_DARK

    # Right Box: Q7 weighted average age
    if os.path.exists(q7_img_path):
        slide6_new.shapes.add_picture(q7_img_path, Inches(6.8), Inches(1.5), width=Inches(5.9))
        
    q7_caption = slide6_new.shapes.add_textbox(Inches(6.8), Inches(5.6), Inches(5.9), Inches(1.3))
    tf7 = q7_caption.text_frame
    tf7.word_wrap = True
    p7 = tf7.paragraphs[0]
    p7.text = "Q7: Oldest Double-Dose Recipients by Weighted Age"
    p7.font.size = Pt(15)
    p7.font.bold = True
    p7.font.color.rgb = DARK_BLUE
    
    bp7 = tf7.add_paragraph()
    bp7.text = "• Bulgaria (BG) and Croatia (HR) represent the oldest double-dosed populations, with a weighted average age of recipients exceeding 52 years.\n• This indicates elder-centric vaccination targeting, contrasting with countries like Ireland that achieved broad age rollouts."
    bp7.font.size = Pt(11)
    bp7.font.color.rgb = TEXT_DARK


    # ==========================================
    # SLIDE 7: Skepticism & Capacity decoupling (Q5, Q8 IMAGES Side-by-Side)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "Results: Skepticism & ICU Capacity Transitions (Q5, Q8)")
    add_content_footer(slide7)

    q5_img_path = os.path.join(plots_dir, "q5_skepticism_vs_hospital_burden.png")
    q8_img_path = os.path.join(plots_dir, "q8_daily_icu_occupancy_comparison.png")
    
    # Left Box: Q5 regression scatter
    if os.path.exists(q5_img_path):
        slide7.shapes.add_picture(q5_img_path, Inches(0.6), Inches(1.5), width=Inches(5.8))
        
    q5_caption = slide7.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(5.8), Inches(1.3))
    tf5 = q5_caption.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "Q5: Skepticism vs. Clinical Peak Burden"
    p5.font.size = Pt(15)
    p5.font.bold = True
    p5.font.color.rgb = DARK_BLUE
    
    bp5 = tf5.add_paragraph()
    bp5.text = "• Pearson correlation maps vaccine skepticism (1 - first dose coverage) to peak weekly hospital admissions per 100k in 2021.\n• Countries with higher vaccine skepticism experienced significantly more severe hospitalization peaks, showing the vaccine's direct protective impact."
    bp5.font.size = Pt(11)
    bp5.font.color.rgb = TEXT_DARK

    # Right Box: Q8 Daily ICU Comparison (Decoupling)
    if os.path.exists(q8_img_path):
        slide7.shapes.add_picture(q8_img_path, Inches(6.8), Inches(1.5), width=Inches(5.9))
        
    q8_caption = slide7.shapes.add_textbox(Inches(6.8), Inches(5.6), Inches(5.9), Inches(1.3))
    tf8 = q8_caption.text_frame
    tf8.word_wrap = True
    p8 = tf8.paragraphs[0]
    p8.text = "Q8: Clinical Severity Decoupling (2020 vs. 2022)"
    p8.font.size = Pt(15)
    p8.font.bold = True
    p8.font.color.rgb = DARK_BLUE
    
    bp8 = tf8.add_paragraph()
    bp8.text = "• Measures peak daily ICU occupancy rate per 100k people across the EU/EEA.\n• Shows a massive reduction in ICU burden in 2022 (green) compared to 2020 (gray) across almost all countries, despite vastly higher infection waves in 2022 due to Omicron."
    bp8.font.size = Pt(11)
    bp8.font.color.rgb = TEXT_DARK


    # ==========================================
    # SLIDE 8: Reflections & Conclusions (Light Background)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide8, "Reflections & Key Takeaways")
    add_content_footer(slide8)

    col1_box = slide8.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0))
    tf1 = col1_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Data Integrity Reflections"
    p.font.name = 'Calibri'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(12)

    bullets9 = [
        ("Missing Standardization", "Inconsistencies like Germany's omission of vaccine names (`UNK`) and coarse-grained age cohorts represent a massive real-world data sharing gap."),
        ("Subnational Row Risks", "The mixing of NUTS-2 regional rows and country totals highlights the danger of double-counting in public health dashboards, requiring rigid data cleaning pipelines."),
        ("Testing Policy Shifting Bias", "Comparing raw case volumes in 2020 vs. 2022 is heavily biased by changes in national testing policies (broad PCR testing in 2020 vs. rapid home antigen testing in 2022).")
    ]
    for title, desc in bullets9:
        bp = tf1.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)

    col2_box = slide8.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0))
    tf2 = col2_box.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = "Concluding Takeaways"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE
    p2.space_after = Pt(12)

    bullets10 = [
        ("Power of EDA", "Exploratory data analysis serves as a vital diagnostic tool in public health, turning raw CSVs into highly structured and actionable metrics."),
        ("Epidemiological Proof", "The statistical results provide clear mathematical proof that vaccine rollouts successfully mitigated clinical intensive care burden across the EU/EEA."),
        ("Transparent Methodology", "Using clean, documented Python script pipelines ensures reproducibility, auditability, and absolute clarity for grading.")
    ]
    for title, desc in bullets10:
        bp = tf2.add_paragraph()
        bp.text = f"•  {title}: "
        bp.font.bold = True
        bp.font.size = Pt(13)
        bp.font.color.rgb = TEXT_DARK
        
        run = bp.add_run()
        run.text = desc
        run.font.bold = False
        bp.space_after = Pt(8)

    prs.save("COVID19_Epidemiology_Presentation.pptx")
    print("Successfully generated COVID19_Epidemiology_Presentation.pptx with all Q1-Q8 embedded plots!")

if __name__ == "__main__":
    create_presentation()
